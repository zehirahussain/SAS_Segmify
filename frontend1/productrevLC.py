import pandas as pd
import matplotlib.pyplot as plt
import os
import mysql.connector
from pptx import Presentation
from pptx.util import Inches
import json

# Database connection function
def get_db_connection():
    return mysql.connector.connect(
        host="localhost",
        user="root",
        password="",
        database="loginandanalysis"
    )

# Function to save image details to the database
def save_or_update_image_in_db(user_id, image_path, image_type):
    conn = get_db_connection()
    cursor = conn.cursor()

    cursor.execute(
        "SELECT id FROM user_images WHERE user_id = %s AND image_path = %s AND image_type = %s",
        (user_id, image_path, image_type)
    )
    result = cursor.fetchone()

    if result:
        cursor.execute(
            "UPDATE user_images SET image_path = %s WHERE id = %s",
            (image_path, result[0])
        )
    else:
        cursor.execute(
            "INSERT INTO user_images (user_id, image_path, image_type) VALUES (%s, %s, %s)",
            (user_id, image_path, image_type)
        )

    conn.commit()
    cursor.close()
    conn.close()

# Function to update PowerPoint presentation with new image
def update_presentation(user_id, image_path):
    conn = get_db_connection()
    cursor = conn.cursor()

    # Fetch or create the presentation for the user
    cursor.execute("SELECT presentation_path FROM user_presentations WHERE user_id = %s", (user_id,))
    result = cursor.fetchone()

    if result:
        ppt_path = result[0]
    else:
        ppt_path = f'static/presentations/user{user_id}_presentation.pptx'
        prs = Presentation()
        prs.save(ppt_path)
        cursor.execute(
            "INSERT INTO user_presentations (user_id, presentation_path) VALUES (%s, %s)",
            (user_id, ppt_path)
        )
        conn.commit()

    prs = Presentation(ppt_path)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.add_picture(image_path, Inches(1), Inches(1), width=Inches(8), height=Inches(5.5))
    
    prs.save(ppt_path)

    cursor.close()
    conn.close()

# Read the dataset from Excel
uploads_dir = "uploads"

# List all files in the uploads directory
files = os.listdir(uploads_dir)

# Identify the Excel file (assuming there is only one Excel file in the directory)
excel_files = [file for file in files if file.endswith('.xlsx')]

# Check if there is at least one Excel file
if not excel_files:
    raise FileNotFoundError("No Excel file found in the uploads directory.")

# Read the first Excel file found
excel_file_path = os.path.join(uploads_dir, excel_files[0])
df = pd.read_excel(excel_file_path)

# Filter relevant columns
relevant_cols = ["Product", "Revenue Billed"]
df = df[relevant_cols]

# Aggregate Revenue by Product
revenue_by_product = df.groupby("Product")["Revenue Billed"].sum().reset_index()

# Bar chart with white background
plt.figure(figsize=(10, 8), facecolor='white')  # Set figure background color to white
plt.bar(revenue_by_product["Product"], revenue_by_product["Revenue Billed"], color='skyblue')
plt.xlabel("Product", fontsize=10, color='black')  # Set x-axis label color to black
plt.ylabel("Total Revenue Billed", fontsize=10, color='black')  # Set y-axis label color to black
plt.title("Total Revenue by Product", fontsize=12, color='black')  # Set title color to black
plt.xticks(rotation=45, ha="right", fontsize=10, color='black')  # Set x-ticks color to black
plt.yticks(fontsize=10, color='black')  # Set y-ticks color to black

# Set background color of the plot area to white
plt.gca().set_facecolor('white')
plt.gca().spines['top'].set_color('black')
plt.gca().spines['right'].set_color('black')
plt.gca().spines['left'].set_color('black')
plt.gca().spines['bottom'].set_color('black')

plt.tight_layout()

# Save the plot as a static image
output_dir = os.path.join('C:\\xampp\\htdocs\\fyp0.3\\frontend1\\static\\images')  # Adjust the path to your XAMPP directory
if not os.path.exists(output_dir):
    os.makedirs(output_dir)
plot_path = os.path.join(output_dir, 'revenue_by_product_bar_chart.png')
plt.savefig(plot_path, facecolor='white')  # Save the figure with white background
plt.close()

# Define user ID and image type
user_id = 1  # Replace with actual user ID
image_type = "Revenue by Product"

# Save image details to the database
save_or_update_image_in_db(user_id, plot_path, image_type)

# Update PowerPoint presentation with the new image
update_presentation(user_id, plot_path)

# Read the analysis results from the file
results_file = 'static/analysis_results.json'
with open(results_file, 'r') as file:
    results = json.load(file)

# Use the analysis result for "revenue_by_product_bar_chart"
analysis_text = results.get('revenue_by_product_bar_chart', 'No analysis result available.')

# Function to append analysis results to the presentation
def append_analysis_to_presentation(user_id, analysis_text):
    conn = get_db_connection()
    cursor = conn.cursor()

    # Fetch the presentation path
    cursor.execute("SELECT presentation_path FROM user_presentations WHERE user_id = %s", (user_id,))
    result = cursor.fetchone()

    if result:
        ppt_path = result[0]
    else:
        print("Presentation not found.")
        cursor.close()
        conn.close()
        return

    prs = Presentation(ppt_path)
    
    # Add analysis text on a new slide
    analysis_slide = prs.slides.add_slide(prs.slide_layouts[5])
    textbox = analysis_slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5.5))
    text_frame = textbox.text_frame
    p = text_frame.add_paragraph()
    p.text = analysis_text
    p.font.size = Inches(0.2)  # Adjust font size as needed

    prs.save(ppt_path)
    print(f"Updated presentation with analysis at {ppt_path}")

    cursor.close()
    conn.close()

# Append the analysis results to the presentation
append_analysis_to_presentation(user_id, analysis_text)
