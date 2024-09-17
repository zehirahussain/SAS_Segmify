import pandas as pd
import matplotlib.pyplot as plt
import os
import json
import mysql.connector
from pptx import Presentation
from pptx.util import Inches





# Database connection function
def get_db_connection():
    return mysql.connector.connect(
        host="localhost",
        user="root",
        password="",
        database="loginandanalysis"
    )

# Function to save or update image details in the database
def save_or_update_image_in_db(user_id, image_path, image_type):
    conn = get_db_connection()
    cursor = conn.cursor()

    cursor.execute(
        "SELECT id FROM user_images WHERE user_id = %s AND image_path = %s AND image_type = %s",
        (user_id, image_path, image_type)
    )
    result = cursor.fetchone()

    if result:
        cursor.execute(
            "UPDATE user_images SET image_path = %s WHERE id = %s",
            (image_path, result[0])
        )
    else:
        cursor.execute(
            "INSERT INTO user_images (user_id, image_path, image_type) VALUES (%s, %s, %s)",
            (user_id, image_path, image_type)
        )

    conn.commit()
    cursor.close()
    conn.close()

# Function to update PowerPoint presentation with new image and analysis results
def update_presentation(user_id, image_path, analysis_text):
    conn = get_db_connection()
    cursor = conn.cursor()

    # Fetch or create the presentation for the user
    cursor.execute("SELECT presentation_path FROM user_presentations WHERE user_id = %s", (user_id,))
    result = cursor.fetchone()

    if result:
        ppt_path = result[0]
    else:
        ppt_path = f'static/presentations/user{user_id}_presentation.pptx'
        prs = Presentation()
        prs.save(ppt_path)
        cursor.execute(
            "INSERT INTO user_presentations (user_id, presentation_path) VALUES (%s, %s)",
            (user_id, ppt_path)
        )
        conn.commit()
# Use absolute path for the presentation file
    ppt_path = 'C:/xampp/htdocs/fyp0.3/static/presentations/user1_presentation.pptx'

    if os.path.exists(ppt_path):
        prs = Presentation(ppt_path)
    # Your code to process the presentation
    else:
       print(f"File not found: {ppt_path}")
    #prs = Presentation(ppt_path)

    # Add the image slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.add_picture(image_path, Inches(1), Inches(1), width=Inches(8), height=Inches(5.5))

    # Add analysis text on a new slide
    analysis_slide = prs.slides.add_slide(prs.slide_layouts[5])
    textbox = analysis_slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5.5))
    text_frame = textbox.text_frame
    p = text_frame.add_paragraph()
    p.text = analysis_text
    p.font.size = Inches(0.2)  # Reduced font size for analysis text

    prs.save(ppt_path)

    cursor.close()
    conn.close()

# Load the dataset
#file_path = os.path.join(os.getcwd(), 'review', 'Review_analysis.xlsx')
file_path = 'C:/xampp/htdocs/fyp0.3/frontend1/review/Review_analysis.xlsx'

df = pd.read_excel(file_path)

#file_path = 'review/Review_analysis.xlsx'
print(os.path.abspath(file_path))
#df = pd.read_excel(file_path)

# Ensure static/images/ exists
output_dir_json = 'C:/xampp/htdocs/fyp0.3/frontend1/review/'
output_dir = 'C:/xampp/htdocs/fyp0.3/frontend1/static/images/'
if not os.path.exists(output_dir):
    os.makedirs(output_dir)

# Semantic analysis storage
semantic_analysis = {
    'revenue_by_item_currency': '',
    'revenue_quantity_by_bu': '',
    'top_customers_by_revenue': '',
}

# Define user ID
user_id = 1  # Replace with actual user ID

# 1. Revenue by Item Name and Currency
revenue_item_currency = df.groupby(['Item Name', 'Currency'])['Revenue Billed'].sum().unstack()
ax = revenue_item_currency.plot(kind='bar', stacked=True, figsize=(20, 12))
plt.title('Revenue by Item Name and Currency')
plt.ylabel('Revenue Billed')
plt.xticks(rotation=90)  # Rotate x-axis labels
plt.tight_layout()  # Adjust layout to prevent clipping
image_path_1 = output_dir + 'revenue_by_item_currency.png'
plt.savefig(image_path_1, bbox_inches='tight', dpi=50)
plt.close()

# Save image details to the database
save_or_update_image_in_db(user_id, image_path_1, 'Revenue by Item and Currency')

# Find top items and currencies with significant revenue
item_totals = revenue_item_currency.sum(axis=1)
top_items = item_totals.nlargest(2).index
top_currencies = {item: revenue_item_currency.loc[item].idxmax() for item in top_items}

top_item_x, top_item_y = top_items[0], top_items[1]
top_currency_x, top_currency_y = top_currencies.get(top_item_x, 'Unknown'), top_currencies.get(top_item_y, 'Unknown')

# Add semantic analysis for revenue by item name and currency
semantic_analysis['revenue_by_item_currency'] = f"Significant revenue contributions are observed for items '{top_item_x}' and '{top_item_y}', primarily in currencies '{top_currency_x}' and '{top_currency_y}'."

# Update PowerPoint presentation with the image and analysis results
#update_presentation(user_id, image_path_1, semantic_analysis['revenue_by_item_currency'])

# 4. Revenue and Quantity by Business Unit (BU)
bu_revenue = df.groupby('BU')['Revenue Billed'].sum()
bu_quantity = df.groupby('BU')['Quantity Billed'].sum()

fig, ax1 = plt.subplots(figsize=(12, 8))
ax1.set_xlabel('Business Unit')
ax1.set_ylabel('Revenue Billed', color='tab:blue')
ax1.bar(bu_revenue.index, bu_revenue, color='tab:blue', label='Revenue Billed')
ax1.tick_params(axis='y', labelcolor='tab:blue')
ax1.set_xticklabels(bu_revenue.index, rotation=90)  # Rotate x-axis labels

ax2 = ax1.twinx()
ax2.set_ylabel('Quantity Billed', color='tab:orange')
ax2.plot(bu_quantity.index, bu_quantity, color='tab:orange', label='Quantity Billed', marker='o')
ax2.tick_params(axis='y', labelcolor='tab:orange')

plt.title('Revenue and Quantity by Business Unit')
plt.tight_layout()  # Adjust layout to prevent clipping
image_path_2 = output_dir + 'revenue_quantity_by_bu.png'
plt.savefig(image_path_2, bbox_inches='tight', dpi=80)
plt.close()

# Save image details to the database
save_or_update_image_in_db(user_id, image_path_2, 'Revenue and Quantity by BU')

# Find top business units for revenue and quantity
top_bu_revenue = bu_revenue.idxmax()
top_bu_quantity = bu_quantity.idxmax()

# Add semantic analysis for revenue and quantity by business unit
semantic_analysis['revenue_quantity_by_bu'] = f"The unit '{top_bu_revenue}' generated the highest revenue, while '{top_bu_quantity}' had the highest quantity billed."

# Update PowerPoint presentation with the image and analysis results
#update_presentation(user_id, image_path_2, semantic_analysis['revenue_quantity_by_bu'])

# 5. Top Customers by Revenue
top_customers = df.groupby('Customer Name')['Revenue Billed'].sum().nlargest(10)
ax = top_customers.plot(kind='bar', figsize=(10, 6))
plt.title('Top 10 Customers by Revenue')
plt.ylabel('Revenue Billed')
plt.xticks(rotation=90)  # Rotate x-axis labels
plt.tight_layout()  # Adjust layout to prevent clipping
image_path_3 = output_dir + 'top_customers_by_revenue.png'
plt.savefig(image_path_3, bbox_inches='tight', dpi=100)
plt.close()

# Save image details to the database
save_or_update_image_in_db(user_id, image_path_3, 'Top Customers by Revenue')

# Find top customer by revenue
top_customer = top_customers.idxmax()
top_customer_revenue = top_customers.max()
semantic_analysis['top_customers_by_revenue'] = f"'{top_customer}' is the top customer with a total revenue of {top_customer_revenue:.2f}."

# Update PowerPoint presentation with the image and analysis results
#update_presentation(user_id, image_path_3, semantic_analysis['top_customers_by_revenue'])

# Save semantic analysis to a JSON file
with open(output_dir_json + 'semantic_analysis.json', 'w') as f:
    json.dump(semantic_analysis, f, indent=4)

print("Review analysis images and semantic analysis have been saved.")

update_presentation(user_id, image_path_1, semantic_analysis['revenue_by_item_currency'])
update_presentation(user_id, image_path_2, semantic_analysis['revenue_quantity_by_bu'])
update_presentation(user_id, image_path_3, semantic_analysis['top_customers_by_revenue'])
