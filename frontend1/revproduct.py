import pandas as pd
import matplotlib.pyplot as plt
import os
import mysql.connector
from pptx import Presentation
from pptx.util import Inches

# Database connection function
def get_db_connection():
    return mysql.connector.connect(
        host="localhost",
        user="root",
        password="",
        database="loginandanalysis"
    )

# Function to save or update image details in the database
def save_or_update_image_in_db(user_id, image_path, image_type):
    conn = get_db_connection()
    cursor = conn.cursor()

    cursor.execute(
        "SELECT id FROM user_images WHERE user_id = %s AND image_path = %s AND image_type = %s",
        (user_id, image_path, image_type)
    )
    result = cursor.fetchone()

    if result:
        cursor.execute(
            "UPDATE user_images SET image_path = %s WHERE id = %s",
            (image_path, result[0])
        )
    else:
        cursor.execute(
            "INSERT INTO user_images (user_id, image_path, image_type) VALUES (%s, %s, %s)",
            (user_id, image_path, image_type)
        )

    conn.commit()
    cursor.close()
    conn.close()

# Function to update PowerPoint presentation with new image and analysis results
def update_presentation(user_id, image_path, analysis_text):
    conn = get_db_connection()
    cursor = conn.cursor()

    # Fetch or create the presentation for the user
    cursor.execute("SELECT presentation_path FROM user_presentations WHERE user_id = %s", (user_id,))
    result = cursor.fetchone()

    if result:
        ppt_path = result[0]
    else:
        ppt_path = f'static/presentations/user{user_id}_presentation.pptx'
        prs = Presentation()
        prs.save(ppt_path)
        cursor.execute(
            "INSERT INTO user_presentations (user_id, presentation_path) VALUES (%s, %s)",
            (user_id, ppt_path)
        )
        conn.commit()

    prs = Presentation(ppt_path)

    # Add the image slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.add_picture(image_path, Inches(1), Inches(1), width=Inches(8), height=Inches(5.5))

    # Add analysis text on a new slide
    analysis_slide = prs.slides.add_slide(prs.slide_layouts[5])
    textbox = analysis_slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5.5))
    text_frame = textbox.text_frame
    p = text_frame.add_paragraph()
    p.text = analysis_text
    p.font.size = Inches(0.2)  # Reduced font size for analysis text

    prs.save(ppt_path)

    cursor.close()
    conn.close()

# Read the dataset from Excel
uploads_dir = "uploads"

# List all files in the uploads directory
files = os.listdir(uploads_dir)

# Identify the Excel file (assuming there is only one Excel file in the directory)
excel_files = [file for file in files if file.endswith('.xlsx')]

# Check if there is at least one Excel file
if not excel_files:
    raise FileNotFoundError("No Excel file found in the uploads directory.")

# Read the first Excel file found
excel_file_path = os.path.join(uploads_dir, excel_files[0])
df = pd.read_excel(excel_file_path)

# Ensure that dates are in datetime format
df['Invoice Date'] = pd.to_datetime(df['Invoice Date'])
df['Billing Date'] = pd.to_datetime(df['Billing Date'])

# Calculate Recency
df['Recency'] = (df['Invoice Date'].max() - df['Invoice Date']).dt.days

# Calculate Frequency
frequency_df = df.groupby('Customer Name').agg({'Invoice': 'nunique'}).reset_index()
frequency_df.rename(columns={'Invoice': 'Frequency'}, inplace=True)

# Calculate Monetary
monetary_df = df.groupby('Customer Name').agg({'Revenue Billed': 'sum'}).reset_index()
monetary_df.rename(columns={'Revenue Billed': 'Monetary'}, inplace=True)

# Merge Recency, Frequency, and Monetary
rfm_df = df.groupby('Customer Name').agg({
    'Recency': 'min',
}).reset_index()

rfm_df = rfm_df.merge(frequency_df, on='Customer Name')
rfm_df = rfm_df.merge(monetary_df, on='Customer Name')

# Assign RFM Scores
rfm_df['R'] = pd.qcut(rfm_df['Recency'], 4, labels=False, duplicates='drop')
rfm_df['F'] = pd.qcut(rfm_df['Frequency'].rank(method="first"), 4, labels=False, duplicates='drop')
rfm_df['M'] = pd.qcut(rfm_df['Monetary'], 4, labels=False, duplicates='drop')

# Create a combined RFM score
rfm_df['RFM_Score'] = rfm_df['R'].astype(str) + rfm_df['F'].astype(str) + rfm_df['M'].astype(str)

# Segment customers based on RFM Score
def segment_customer(df):
    if df['RFM_Score'] == '333':
        return 'Best Customers'
    elif df['RFM_Score'].startswith('33'):
        return 'Loyal Customers'
    elif df['RFM_Score'].startswith('3'):
        return 'Potential Loyalists'
    elif df['RFM_Score'].startswith('2'):
        return 'At Risk'
    else:
        return 'Others'

rfm_df['Segment'] = rfm_df.apply(segment_customer, axis=1)

# Pie chart for customer segmentation
segment_counts = rfm_df['Segment'].value_counts()

plt.figure(figsize=(10, 8))
plt.pie(segment_counts, labels=segment_counts.index, autopct='%1.1f%%', startangle=140, 
        colors=plt.cm.Paired(range(len(segment_counts))), textprops={'fontsize': 12, 'color': 'black'})
plt.title('Customer Segmentation using RFM', fontsize=14, color='black')
plt.axis('equal')  # Equal aspect ratio ensures that pie is drawn as a circle.

# Save the plot as a static image
output_dir = os.path.join('C:\\xampp\\htdocs\\fyp0.3\\frontend1\\static\\images')
if not os.path.exists(output_dir):
    os.makedirs(output_dir)
plot_path = os.path.join(output_dir, 'revenue_by_product_pie_chart.png')
plt.savefig(plot_path, facecolor='white')  # Save the figure with white background
plt.close()

print(f'Pie chart for customer segmentation using RFM saved to {plot_path}')

# Define user ID and image type
user_id = 1  # Replace with actual user ID
image_type = "Revenue by Product"

# Save or update image details in the database
save_or_update_image_in_db(user_id, plot_path, image_type)

# Example analysis text
analysis_text = "This pie chart shows the distribution of customers based on their RFM scores. Each segment represents a different customer group such as Best Customers, Loyal Customers, Potential Loyalists, and At Risk."

# Update PowerPoint presentation with the new image and analysis results
update_presentation(user_id, plot_path, analysis_text)
