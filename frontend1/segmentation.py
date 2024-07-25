import pandas as pd
import matplotlib.pyplot as plt
import os
import mysql.connector
from pptx import Presentation
from pptx.util import Inches

# Database connection function
def get_db_connection():
    return mysql.connector.connect(
        host="localhost",
        user="root",
        password="",
        database="loginandanalysis"
    )

# Function to save image details to the database
def save_image_to_db(user_id, image_path, image_type):
    conn = get_db_connection()
    cursor = conn.cursor()
    cursor.execute(
        "INSERT INTO user_images (user_id, image_path, image_type) VALUES (%s, %s, %s)",
        (user_id, image_path, image_type)
    )
    conn.commit()
    cursor.close()
    conn.close()

# Function to update PowerPoint presentation with new image
def update_presentation(user_id, image_path):
    conn = get_db_connection()
    cursor = conn.cursor()

    # Fetch or create the presentation for the user
    cursor.execute("SELECT presentation_path FROM user_presentations WHERE user_id = %s", (user_id,))
    result = cursor.fetchone()

    if result:
        ppt_path = result[0]
    else:
        ppt_path = f'static/presentations/user{user_id}_presentation.pptx'
        prs = Presentation()
        prs.save(ppt_path)
        cursor.execute(
            "INSERT INTO user_presentations (user_id, presentation_path) VALUES (%s, %s)",
            (user_id, ppt_path)
        )
        conn.commit()

    prs = Presentation(ppt_path)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.add_picture(image_path, Inches(1), Inches(1), width=Inches(8), height=Inches(5.5))
    prs.save(ppt_path)

    cursor.close()
    conn.close()

# Read the dataset from Excel
uploads_dir = "uploads"

# List all files in the uploads directory
files = os.listdir(uploads_dir)

# Identify the Excel file (assuming there is only one Excel file in the directory)
excel_files = [file for file in files if file.endswith('.xlsx')]

# Check if there is at least one Excel file
if not excel_files:
    raise FileNotFoundError("No Excel file found in the uploads directory.")

# Read the first Excel file found
excel_file_path = os.path.join(uploads_dir, excel_files[0])
df = pd.read_excel(excel_file_path)

# Filter relevant columns
relevant_cols = ["BU", "Revenue Billed", "Currency", "Quantity Billed", "SB FX Rate", "Sales Tax Billed", "Recog. Months"]
df = df[relevant_cols]

# Aggregate MRR by Business Unit
mrr_by_bu = df.groupby("BU")["Revenue Billed"].sum().reset_index()

# print on Visula studio code console
print("Aggregated Revenue by BU:")
print(mrr_by_bu)



# Bar chart with white background
plt.figure(figsize=(10, 8), facecolor='white')  # Set figure background color to white
plt.bar(mrr_by_bu["BU"], mrr_by_bu["Revenue Billed"], color='skyblue')

# Set axis labels and title with black color
plt.xlabel("Business Unit", fontsize=10, color='black')
plt.ylabel("Monthly Recurring Revenue (MRR)", fontsize=10, color='black')
plt.title("Monthly Recurring Revenue by Business Unit", fontsize=12, color='black')

# Set x and y ticks with black color
plt.xticks(rotation=80, ha="right", fontsize=8, color='black')
plt.yticks(fontsize=10, color='black')

# Set background color of the plot area to white
plt.gca().set_facecolor('white')
plt.gca().spines['top'].set_color('black')
plt.gca().spines['right'].set_color('black')
plt.gca().spines['left'].set_color('black')
plt.gca().spines['bottom'].set_color('black')

# Add grid lines with black color for better readability
plt.grid(color='black', linestyle='--', linewidth=0.5)

plt.tight_layout()

# Save the plot as a static image
output_dir = os.path.join('C:\\xampp\\htdocs\\fyp0.3\\frontend1\\static\\images')  # Adjust the path to your XAMPP directory
if not os.path.exists(output_dir):
    os.makedirs(output_dir)
plot_path = os.path.join(output_dir, 'mrr_by_bu.png')
plt.savefig(plot_path, facecolor='white')  # Save the figure with white background
plt.close()

# Define user ID and image type
user_id = 1  # Replace with actual user ID
image_type = "MRR by BU"

# Save image details to the database
save_image_to_db(user_id, plot_path, image_type)

# Update PowerPoint presentation with the new image
update_presentation(user_id, plot_path)
